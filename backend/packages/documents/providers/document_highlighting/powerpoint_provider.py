from __future__ import annotations

import asyncio
from typing import List
from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor

from .interface import DocumentHighlightingInterface, HighlightData
from common.core.otel_axiom_exporter import trace_span, get_logger

logger = get_logger(__name__)


class PowerPointHighlightProvider(DocumentHighlightingInterface):
    """PowerPoint-specific provider for highlighting text using python-pptx."""

    def __init__(self):
        self.supported_extensions = [".pptx", ".ppt"]

    @trace_span
    async def highlight_document(
        self,
        document_content: bytes,
        document_filename: str,
        highlights: List[HighlightData],
    ) -> bytes:
        """Highlight text in a PowerPoint document using python-pptx."""
        if not self.supports_file_type(document_filename):
            raise NotImplementedError(f"File type not supported: {document_filename}")

        if not highlights:
            return document_content

        # Run the synchronous PowerPoint processing in a thread pool
        return await asyncio.to_thread(
            self._highlight_powerpoint_sync, document_content, highlights
        )

    def _highlight_powerpoint_sync(
        self, document_content: bytes, highlights: List[HighlightData]
    ) -> bytes:
        """Synchronous PowerPoint highlighting logic."""
        try:
            # Open PowerPoint from bytes
            input_stream = BytesIO(document_content)
            prs = Presentation(input_stream)
            input_stream.close()

            # Track processed highlights to avoid duplicates
            processed_texts = set()
            total_highlights_applied = 0

            for highlight_data in highlights:
                text_to_find = highlight_data.text.strip()

                # Skip if we've already processed this exact text
                if text_to_find in processed_texts:
                    continue
                processed_texts.add(text_to_find)

                logger.info(
                    f"Searching for text to highlight: '{text_to_find[:50]}...'"
                )

                # Search for text across all slides
                found_instances = 0
                for slide_num, slide in enumerate(prs.slides):
                    # Check all shapes on the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            # Search within text frame
                            found_instances += self._highlight_text_in_text_frame(
                                shape.text_frame, text_to_find, highlight_data
                            )

                if found_instances == 0:
                    logger.warning(
                        f"Text not found in presentation: '{text_to_find[:50]}...'"
                    )
                else:
                    logger.info(
                        f"Highlighted {found_instances} instances of text: '{text_to_find[:50]}...'"
                    )
                    total_highlights_applied += found_instances

            # Save the modified PowerPoint to bytes
            output_stream = BytesIO()
            prs.save(output_stream)
            highlighted_content = output_stream.getvalue()
            output_stream.close()

            logger.info(
                f"Successfully highlighted PowerPoint with {len(highlights)} highlight requests, {total_highlights_applied} instances highlighted"
            )
            return highlighted_content

        except Exception as e:
            logger.error(f"Error highlighting PowerPoint: {e}")
            raise

    def _highlight_text_in_text_frame(
        self, text_frame, text_to_find: str, highlight_data: HighlightData
    ) -> int:
        """
        Highlight text within a text frame.

        Returns:
            Number of instances highlighted
        """
        found_instances = 0

        for paragraph in text_frame.paragraphs:
            # Get the full paragraph text
            paragraph_text = paragraph.text

            if text_to_find.lower() in paragraph_text.lower():
                # Find all occurrences in this paragraph
                start_pos = 0
                while True:
                    pos = paragraph_text.lower().find(text_to_find.lower(), start_pos)
                    if pos == -1:
                        break

                    # Try to highlight the text by modifying runs
                    self._highlight_text_in_paragraph(
                        paragraph, pos, pos + len(text_to_find), highlight_data
                    )
                    found_instances += 1
                    start_pos = pos + 1

        return found_instances

    def _highlight_text_in_paragraph(
        self, paragraph, start_char: int, end_char: int, highlight_data: HighlightData
    ):
        """
        Highlight specific character range in a paragraph.
        Note: PowerPoint highlighting is limited - we'll change the text color as a fallback.
        """
        try:
            # Convert RGB (0-1) to RGB (0-255)
            rgb_color = RGBColor(
                int(highlight_data.color[0] * 255),
                int(highlight_data.color[1] * 255),
                int(highlight_data.color[2] * 255),
            )

            # Find the runs that contain our target text
            current_pos = 0
            for run in paragraph.runs:
                run_length = len(run.text)
                run_start = current_pos
                run_end = current_pos + run_length

                # Check if this run overlaps with our target range
                if run_start < end_char and run_end > start_char:
                    # Apply highlighting (background color) if supported, otherwise text color
                    if hasattr(run.font, "highlight_color"):
                        # Try to set highlight color
                        run.font.highlight_color.rgb = rgb_color
                    else:
                        # Fallback to changing text color
                        run.font.color.rgb = rgb_color

                current_pos = run_end

        except Exception as e:
            logger.warning(f"Could not apply highlighting to paragraph: {e}")

    def supports_file_type(self, filename: str) -> bool:
        """Check if the file is a PowerPoint document."""
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)

    def get_supported_extensions(self) -> List[str]:
        """Return list of supported file extensions."""
        return self.supported_extensions.copy()
